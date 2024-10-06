import argparse
import json
import boto3
import os
from botocore.exceptions import ClientError
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_bedrock_client():
    region = os.environ.get('AWS_REGION')
    if not region:
        raise ValueError("AWS_REGION environment variable is not set")
    return boto3.client(
        service_name='bedrock-runtime',
        region_name=region
    )

def translate_text(text, bedrock_client, source_language_code, target_language_code):
    try:
        prompt = f"""Translate the following text from {source_language_code} to {target_language_code}. 
Provide only the translated text without any explanations or additional text.

Text to translate: {text}"""
        
        body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 1000,
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "temperature": 0.0
        }
        
        response = bedrock_client.invoke_model(
            body=json.dumps(body),
            modelId="anthropic.claude-3-sonnet-20240229-v1:0",
            contentType="application/json"
        )
        
        response_body = json.loads(response['body'].read())
        translated_text = response_body['content'][0]['text']
        return translated_text.strip()
    except Exception as e:
        print(f'Translation error: {str(e)}')
        return text

def copy_cell_format(source_cell, target_cell):
    # Copy font
    target_cell.font = Font(
        name=source_cell.font.name,
        size=source_cell.font.size,
        bold=source_cell.font.bold,
        italic=source_cell.font.italic,
        color=source_cell.font.color
    )
    
    # Copy alignment
    target_cell.alignment = Alignment(
        horizontal=source_cell.alignment.horizontal,
        vertical=source_cell.alignment.vertical,
        wrap_text=source_cell.alignment.wrap_text
    )
    
    # Copy fill
    if source_cell.fill.fill_type:
        target_cell.fill = PatternFill(
            fill_type=source_cell.fill.fill_type,
            fgColor=source_cell.fill.fgColor,
            bgColor=source_cell.fill.bgColor
        )
    
    # Copy border
    if source_cell.border:
        target_cell.border = Border(
            left=Side(border_style=source_cell.border.left.style, color=source_cell.border.left.color),
            right=Side(border_style=source_cell.border.right.style, color=source_cell.border.right.color),
            top=Side(border_style=source_cell.border.top.style, color=source_cell.border.top.color),
            bottom=Side(border_style=source_cell.border.bottom.style, color=source_cell.border.bottom.color)
        )

def translate_workbook(input_path, output_path, source_language_code, target_language_code, bedrock_client):
    print(f'Translating Excel file: {input_path}')
    original_wb = load_workbook(input_path, data_only=True)
    new_wb = Workbook()
    new_wb.remove(new_wb.active)
    
    for sheet_name in original_wb.sheetnames:
        print(f'Translating sheet: {sheet_name}')
        original_sheet = original_wb[sheet_name]
        new_sheet = new_wb.create_sheet(title=sheet_name)
        
        for column in original_sheet.column_dimensions:
            new_sheet.column_dimensions[column] = original_sheet.column_dimensions[column]
        
        for row_index, row in enumerate(original_sheet.iter_rows(), 1):
            for col_index, cell in enumerate(row, 1):
                new_cell = new_sheet.cell(row=row_index, column=col_index)
                copy_cell_format(cell, new_cell)
                
                if isinstance(cell.value, str) and cell.value.strip():
                    new_cell.value = translate_text(
                        cell.value,
                        bedrock_client,
                        source_language_code,
                        target_language_code
                    )
                else:
                    new_cell.value = cell.value
    
    new_wb.save(output_path)

def translate_shape_text(shape, bedrock_client, source_language_code, target_language_code):
    if hasattr(shape, 'text') and shape.text.strip():
        shape.text = translate_text(
            shape.text,
            bedrock_client,
            source_language_code,
            target_language_code
        )
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            translate_shape_text(subshape, bedrock_client, source_language_code, target_language_code)
    
    if hasattr(shape, 'table'):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cell.text = translate_text(
                        cell.text,
                        bedrock_client,
                        source_language_code,
                        target_language_code
                    )

def translate_presentation(input_path, output_path, source_language_code, target_language_code, bedrock_client):
    print(f'Translating PowerPoint file: {input_path}')
    prs = Presentation(input_path)
    
    for i, slide in enumerate(prs.slides, 1):
        print(f'Translating slide {i} of {len(prs.slides)}')
        for shape in slide.shapes:
            translate_shape_text(shape, bedrock_client, source_language_code, target_language_code)
    
    prs.save(output_path)

def translate_document(input_path, source_language_code, target_language_code, bedrock_client):
    file_extension = os.path.splitext(input_path)[1].lower()
    output_path = input_path.replace(file_extension, f'-{target_language_code}{file_extension}')
    
    if file_extension == '.xlsx':
        translate_workbook(input_path, output_path, source_language_code, target_language_code, bedrock_client)
    elif file_extension == '.pptx':
        translate_presentation(input_path, output_path, source_language_code, target_language_code, bedrock_client)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}. Supported types are .xlsx and .pptx")
    
    return output_path

def main():
    argument_parser = argparse.ArgumentParser(
        description='Translates Excel and PowerPoint files using Amazon Bedrock with Claude 3.5'
    )
    argument_parser.add_argument('source_language_code', type=str, help='Source language code (e.g., en)')
    argument_parser.add_argument('target_language_code', type=str, help='Target language code (e.g., es)')
    argument_parser.add_argument('input_file_path', type=str, help='Path to the input file (.xlsx or .pptx)')
    args = argument_parser.parse_args()

    try:
        bedrock_client = get_bedrock_client()
    except ValueError as e:
        print(f"Error: {e}")
        return

    try:
        output_path = translate_document(
            args.input_file_path,
            args.source_language_code,
            args.target_language_code,
            bedrock_client
        )
        print(f'Translation completed. Output saved to {output_path}')
    except ValueError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")

if __name__ == '__main__':
    main()